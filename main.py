import os
import json
import uuid
import hashlib
import subprocess
from datetime import datetime
from typing import Dict, Any, List

import fitz  # PyMuPDF
from pptx import Presentation
from fastapi import FastAPI, Request, UploadFile, File, Form
from fastapi.responses import HTMLResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from openai import OpenAI

# ----------------- CONFIG -----------------
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_DIR = os.path.join(BASE_DIR, "upload")
AUDIO_DIR = os.path.join(UPLOAD_DIR, "audio")
DOC_DIR = os.path.join(UPLOAD_DIR, "docs")
SESSIONS_PATH = os.path.join(BASE_DIR, "sessions.json")

os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(AUDIO_DIR, exist_ok=True)
os.makedirs(DOC_DIR, exist_ok=True)

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

app = FastAPI(title="Helektron Study Assistant")
app.mount("/static", StaticFiles(directory=os.path.join(BASE_DIR, "static")), name="static")
templates = Jinja2Templates(directory=os.path.join(BASE_DIR, "templates"))

# ----------------- SESSION STORAGE -----------------
def load_sessions() -> Dict[str, Any]:
    if not os.path.exists(SESSIONS_PATH):
        return {}
    with open(SESSIONS_PATH, "r") as f:
        try:
            return json.load(f)
        except json.JSONDecodeError:
            return {}

def save_sessions(sessions: Dict[str, Any]) -> None:
    with open(SESSIONS_PATH, "w") as f:
        json.dump(sessions, f, indent=2)

sessions: Dict[str, Any] = load_sessions()

def get_or_create_session(session_id: str | None) -> str:
    global sessions
    if session_id and session_id in sessions:
        return session_id
    new_id = str(uuid.uuid4())
    sessions[new_id] = {
        "created_at": datetime.utcnow().isoformat(),
        "files": [],   # list of {"name","type","text"}
    }
    save_sessions(sessions)
    return new_id

def add_text_to_session(session_id: str, filename: str, file_type: str, text: str) -> None:
    global sessions
    if session_id not in sessions:
        session_id = get_or_create_session(None)
    sessions[session_id]["files"].append({
        "name": filename,
        "type": file_type,
        "text": text,
        "added_at": datetime.utcnow().isoformat()
    })
    save_sessions(sessions)

def get_session_text(session_id: str) -> str:
    session = sessions.get(session_id)
    if not session:
        return ""
    blocks = [f"--- {f['name']} ({f['type']}) ---\n{f['text']}" for f in session["files"] if f.get("text")]
    return "\n\n".join(blocks)

# ----------------- HELPERS: FILE TEXT EXTRACTION -----------------
def extract_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_pdf(path: str) -> str:
    doc = fitz.open(path)
    texts = []
    for page in doc:
        texts.append(page.get_text())
    doc.close()
    return "\n".join(texts)

def extract_pptx(path: str) -> str:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def transcribe_audio(path: str) -> str:
    # Find whisper binary
    candidates = [
        os.path.join(BASE_DIR, "whisper.cpp", "main"),
        os.path.join(BASE_DIR, "whisper.cpp", "build", "bin", "main"),
        os.path.join(BASE_DIR, "whisper-main", "build", "bin", "whisper-cli"),
    ]
    whisper_bin = next((c for c in candidates if os.path.exists(c)), None)
    if not whisper_bin:
        raise RuntimeError("Whisper binary not found. Build whisper.cpp first.")

    wav_path = path + ".wav"
    subprocess.run(
        ["ffmpeg", "-y", "-i", path, "-acodec", "pcm_s16le", "-ar", "16000", "-ac", "1", wav_path],
        check=True
    )

    out_prefix = path + "_out"
    subprocess.run(
        [whisper_bin, "-m", os.path.join(BASE_DIR, "whisper.cpp", "models", "ggml-base.en.bin"),
         "-f", wav_path, "--output-txt", "--output-file", out_prefix],
        check=True
    )

    txt_path = out_prefix + ".txt"
    with open(txt_path, "r", encoding="utf-8", errors="ignore") as f:
        transcript = f.read()

    # cleanup
    for p in [wav_path, txt_path]:
        if os.path.exists(p):
            os.remove(p)

    return transcript

# ----------------- GPT PROMPTS -----------------
def get_summary_prompt(transcript: str) -> str:
    return f"""
Based on the following combined study materials (lectures, slides, PDFs, notes, transcripts),
create a **structured, detailed, and academically accurate study summary**.

Organize the output into the following clearly labeled sections:

- **Overview**: Briefly describe the overall topic and goals.
- **Key Concepts Introduced**: List and explain major ideas, theories, formulas, or processes.
- **Detailed Topic Breakdown**: Group related ideas and summarize explanations, reasoning steps, and relationships.
- **Important Definitions**: Provide concise definitions for technical terms or domain-specific vocabulary.
- **Examples or Applications**: Summarize any examples, demonstrations, or real-world applications.
- **Main Takeaways**: What a student should remember after studying this material.

Expectations:
- Use clear, concise bullet points.
- Prioritize conceptual correctness.
- Use simple academic language suitable for undergraduate study.
- If information is unclear or missing, label it as `TBD`.

---BEGIN MATERIAL---
{transcript}
---END MATERIAL---
"""

def get_keyterms_prompt(text: str) -> str:
    return f"""
From the combined study materials below, extract **10–20 key terms** with short definitions.

For each term, include:
- The term
- A 1–2 sentence definition
- (Optional) A quick note on why it matters in this context.

Format as a bulleted list.

---BEGIN MATERIAL---
{text}
---END MATERIAL---
"""

def get_questions_prompt(text: str) -> str:
    return f"""
Create **8–12 practice questions** based on the combined study materials below.

Include a mix of:
- Conceptual understanding questions
- Short-answer questions
- Application/problem-style questions (where possible)

Do NOT provide answers, only questions.
Group them into sections if appropriate (e.g., 'Conceptual', 'Application').

---BEGIN MATERIAL---
{text}
---END MATERIAL---
"""

def get_resources_prompt(text: str) -> str:
    return f"""
Based on the topics and concepts in the combined study materials below,
recommend **3–7 high-quality external resources** for further study.

Each resource should include:
- Title
- Type (e.g., textbook chapter, .edu article, video lecture)
- Source (e.g., university, well-known platform)
- 1–2 sentences on why it is relevant.

Prefer:
- .edu domains
- Reputable textbooks
- Well-known educational channels

You do NOT need to provide URLs, just clear, identifiable references.

---BEGIN MATERIAL---
{text}
---END MATERIAL---
"""

def call_gpt(prompt: str) -> str:
    resp = client.chat.completions.create(
        model="gpt-4.1",
        messages=[{"role": "user", "content": prompt}],
    )
    return resp.choices[0].message.content.strip()

# ----------------- ROUTES -----------------
@app.get("/", response_class=HTMLResponse)
def index(request: Request):
    return templates.TemplateResponse("index.html", {"request": request})

# Upload any material (txt, pdf, pptx, mp4, m4a)
@app.post("/upload", response_class=HTMLResponse)
async def upload_material(
    request: Request,
    file: UploadFile = File(...),
    session_id: str | None = Form(None),
):
    global sessions

    session_id = get_or_create_session(session_id)
    filename = file.filename
    ext = filename.split(".")[-1].lower()

    # Save file
    if ext in ["mp4", "m4a", "wav", "webm"]:
        save_dir = AUDIO_DIR
        kind = "audio"
    else:
        save_dir = DOC_DIR
        kind = "document"

    os.makedirs(save_dir, exist_ok=True)
    saved_path = os.path.join(save_dir, f"{uuid.uuid4().hex}_{filename}")

    with open(saved_path, "wb") as f:
        f.write(await file.read())

    # Extract text
    try:
        if ext == "txt":
            text = extract_txt(saved_path)
        elif ext == "pdf":
            text = extract_pdf(saved_path)
        elif ext in ["pptx"]:
            text = extract_pptx(saved_path)
        elif ext in ["mp4", "m4a", "wav", "webm"]:
            text = transcribe_audio(saved_path)
        else:
            text = f"[Unsupported file type: {ext}]"
    except Exception as e:
        text = f"[Error processing file {filename}: {e}]"

    add_text_to_session(session_id, filename, kind, text)

    # Render updated materials list (fragment)
    session = sessions[session_id]
    return templates.TemplateResponse(
        "fragments/upload_status.html",
        {
            "request": request,
            "session_id": session_id,
            "files": session["files"],
        },
    )

# Live transcription: audio blob from the browser
@app.post("/upload_live_audio", response_class=HTMLResponse)
async def upload_live_audio(
    request: Request,
    audio: UploadFile = File(...),
    session_id: str | None = Form(None),
):
    session_id = get_or_create_session(session_id)
    filename = audio.filename or "live_recording.webm"
    ext = filename.split(".")[-1].lower()

    os.makedirs(AUDIO_DIR, exist_ok=True)
    saved_path = os.path.join(AUDIO_DIR, f"live_{uuid.uuid4().hex}.{ext}")

    with open(saved_path, "wb") as f:
        f.write(await audio.read())

    try:
        text = transcribe_audio(saved_path)
    except Exception as e:
        text = f"[Error transcribing live audio: {e}]"

    add_text_to_session(session_id, filename, "audio-live", text)
    session = sessions[session_id]

    return templates.TemplateResponse(
        "fragments/upload_status.html",
        {
            "request": request,
            "session_id": session_id,
            "files": session["files"],
        },
    )

# ---- Study tools (summary, key terms, questions, resources) ----

@app.get("/summary/{session_id}", response_class=HTMLResponse)
def generate_summary(request: Request, session_id: str):
    text = get_session_text(session_id)
    if not text.strip():
        content = "No material uploaded yet."
    else:
        prompt = get_summary_prompt(text)
        content = call_gpt(prompt)

    return templates.TemplateResponse(
        "fragments/summary.html",
        {"request": request, "content": content},
    )

@app.get("/keyterms/{session_id}", response_class=HTMLResponse)
def generate_keyterms(request: Request, session_id: str):
    text = get_session_text(session_id)
    if not text.strip():
        content = "No material uploaded yet."
    else:
        prompt = get_keyterms_prompt(text)
        content = call_gpt(prompt)

    return templates.TemplateResponse(
        "fragments/keyterms.html",
        {"request": request, "content": content},
    )

@app.get("/questions/{session_id}", response_class=HTMLResponse)
def generate_questions_view(request: Request, session_id: str):
    text = get_session_text(session_id)
    if not text.strip():
        content = "No material uploaded yet."
    else:
        prompt = get_questions_prompt(text)
        content = call_gpt(prompt)

    return templates.TemplateResponse(
        "fragments/questions.html",
        {"request": request, "content": content},
    )

@app.get("/resources/{session_id}", response_class=HTMLResponse)
def generate_resources_view(request: Request, session_id: str):
    text = get_session_text(session_id)
    if not text.strip():
        content = "No material uploaded yet."
    else:
        prompt = get_resources_prompt(text)
        content = call_gpt(prompt)

    return templates.TemplateResponse(
        "fragments/resources.html",
        {"request": request, "content": content},
    )

